"""TPS Deck Builder — PowerPoint briefing deck via python-pptx.

Produces a structured intelligence briefing deck:
  Slide 1: Title + BLUF
  Slide 2: Key Judgments
  Slides 3-N: One per section with embedded assets
  Final slide: Sources + QC watermark
"""
from __future__ import annotations

import os
from typing import Optional

from core.schemas import (
    PresentationPlan,
    AssetProvenance,
    DeliverableKind,
)
from core.style_director import get_brand
from core.exceptions import GeneratorError
from deliverables.base import BaseDeliverableBuilder


class DeckBuilder(BaseDeliverableBuilder):
    """Build PowerPoint briefing deck."""

    @property
    def kind(self) -> DeliverableKind:
        return DeliverableKind.DECK

    def build(
        self,
        brief_json: dict,
        plan: PresentationPlan,
        asset_outputs: dict[str, str],
        provenance_records: list[AssetProvenance],
        output_path: str,
    ) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise GeneratorError(
                "python-pptx not installed. Run: pip install python-pptx"
            )

        brand = get_brand()
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        title = brief_json.get("title", "Intelligence Brief")
        bluf = brief_json.get("bluf", "")
        judgments = brief_json.get("headline_judgments", [])
        sections = brief_json.get("sections", [])

        primary_rgb = self._hex_to_rgb(brand.color_primary)
        accent_rgb = self._hex_to_rgb(brand.color_accent)
        body_rgb = self._hex_to_rgb(brand.color_body)

        # ── Slide 1: Title ──────────────────────────────────────────────
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*primary_rgb)

        # Title text
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(13), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True

        # BLUF
        if bluf:
            txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(13), Inches(2))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = bluf[:300]
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(200, 200, 200)

        # QC stamp
        self._add_qc_footer(slide)

        # ── Slide 2: Key Judgments ──────────────────────────────────────
        if judgments:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = RGBColor(245, 243, 239)

            # Header
            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "KEY JUDGMENTS"
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(*primary_rgb)
            p.font.bold = True

            # Judgment list
            y = 1.8
            for j in judgments[:6]:
                claim = j.get("claim", j.get("judgment", ""))
                kent = j.get("kent_band", "")
                band_label = kent.replace("_", " ").title() if kent else ""

                txBox = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(14), Inches(0.9))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                if band_label:
                    run = p.add_run()
                    run.text = f"[{band_label}] "
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(*accent_rgb)
                    run.font.bold = True
                run2 = p.add_run()
                run2.text = claim[:200]
                run2.font.size = Pt(14)
                run2.font.color.rgb = RGBColor(*body_rgb)

                y += 1.0

            self._add_qc_footer(slide)

        # ── Slides 3-N: Sections with assets ───────────────────────────
        asset_list = list(asset_outputs.items())
        asset_idx = 0

        for section in sections[:8]:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = RGBColor(245, 243, 239)

            # Section title
            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = section.get("title", "Section")
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(*primary_rgb)
            p.font.bold = True

            # Narrative text (left half)
            narrative = section.get("narrative", "")
            if narrative:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(7), Inches(6.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = narrative[:500]
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(*body_rgb)

            # Embed asset image (right half)
            if asset_idx < len(asset_list):
                _, asset_path = asset_list[asset_idx]
                if asset_path.endswith(('.png', '.jpg', '.jpeg')):
                    try:
                        slide.shapes.add_picture(
                            asset_path, Inches(8.5), Inches(1.5),
                            width=Inches(6.5), height=Inches(6),
                        )
                    except Exception:
                        pass
                asset_idx += 1

            self._add_qc_footer(slide)

        # ── Final slide: Sources + Provenance ──────────────────────────
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*primary_rgb)

        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1), Inches(13), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "SOURCES & PROVENANCE"
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True

        # Provenance summary
        if provenance_records:
            y = 2.5
            for prov in provenance_records[:8]:
                txBox = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(13), Inches(0.5))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = (
                    f"{prov.asset_id} | {prov.generator} | "
                    f"{prov.model_used} | ${prov.actual_cost_usd:.4f} | "
                    f"{prov.qc_status.value}"
                )
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(180, 180, 180)
                y += 0.5

        self._add_qc_footer(slide, white=True)

        # Save
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)
        return output_path

    def _add_qc_footer(self, slide, white=False):
        """Add QC watermark footer to a slide."""
        try:
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
        except ImportError:
            return

        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(8.3), Inches(14), Inches(0.5)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "PENDING HUMAN ANALYST QC REVIEW | Trevor Intelligence"
        p.font.size = Pt(8)
        if white:
            p.font.color.rgb = RGBColor(120, 120, 120)
        else:
            p.font.color.rgb = RGBColor(150, 150, 150)

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
        """Convert hex color to RGB tuple."""
        h = hex_color.lstrip("#")
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
